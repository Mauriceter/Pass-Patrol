import os
import re
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
from odf.opendocument import load as load_odt
from odf.text import P
import xlrd  # For handling .xls files
import argparse
import datetime
import warnings
import shutil
import json
from concurrent.futures import ThreadPoolExecutor, as_completed
import ollama


# ANSI escape code 
RED_TEXT = '\033[91m'
PURPLE_TEXT = '\033[35m'
BOLD_TEXT = "\033[1m"
RESET_TEXT = '\033[0m'

# Define output folder
timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
download_folder = f"downloaded_files_{timestamp}"
os.makedirs(download_folder, exist_ok=True)

# Initialize the JSON report dictionary
json_report = {}

# Global variable for debug mode
debug = False
interesting_filenames = []

warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")


# Function to load keywords from a specified file
def load_keywords(keywords_file):
    """Load keywords from the specified file."""
    try:
        with open(keywords_file, 'r') as kf:
            return [line.strip().lower() for line in kf.readlines()]
    except Exception as e:
        print(f"Error loading keywords: {e}")
        return []


def search_in_text(text, keywords):
    """Search for keywords in the text, return a list of tuples for all matches."""
    matches = []
    for keyword in keywords:
        keyword_pos = text.lower().find(keyword.lower())
        while keyword_pos != -1:
            matches.append((keyword, keyword_pos))
            keyword_pos = text.lower().find(keyword.lower(), keyword_pos + len(keyword))
    return matches


def create_snippet(text, keyword_tuple):
    """Generate a snippet around a keyword with up to 50 characters and at most 1 line before and 1 line after,
       highlighting the keyword in red."""
    
    keyword = keyword_tuple[0]
    keyword_pos = keyword_tuple[1]

    # Define the character limits around the keyword
    start_char = max(keyword_pos - 50, 0)
    end_char = min(keyword_pos + len(keyword) + 50, len(text))
    context_text = text[start_char:end_char]
    
    # Limit context_text to 1 line before and 1 line after the keyword line
    lines = context_text.splitlines()
    keyword_line_index = next((i for i, line in enumerate(lines) if keyword.lower() in line.lower()), None)
    if keyword_line_index is None:
        return None

    # Extract up to 1 line before and 1 line after the keyword line
    start_line = max(keyword_line_index - 1, 0)
    end_line = min(keyword_line_index + 2, len(lines))  # +2 to include the keyword line and 1 line after
    snippet = '\n'.join(lines[start_line:end_line])

    # Highlight the keyword in red
    highlighted_snippet = re.sub(re.escape(keyword), f"{RED_TEXT}\\g<0>{RESET_TEXT}", snippet, flags=re.IGNORECASE)
    
    return highlighted_snippet


def copy_file_to_download_folder(file_path):
    """Copy file to the downloaded folder and avoid overwriting files with the same name."""
    base_name = os.path.basename(file_path)
    dest_path = os.path.join(download_folder, base_name)
    
    # Check if the file already exists in the download folder and rename it if necessary
    counter = 1
    while os.path.exists(dest_path):
        name, ext = os.path.splitext(base_name)
        dest_path = os.path.join(download_folder, f"{name}_{counter}{ext}")
        counter += 1
    
    shutil.copy(file_path, dest_path)
    return dest_path


# Function to handle PDF files
def handle_pdf_file(file_path, keywords, limit=10):
    """Read and search in PDF files."""
    snippets = []
    try:
        with pdfplumber.open(file_path) as pdf:
            match_count = 0  # Counter for number of matches
            for page in pdf.pages:
                text = page.extract_text()
                matches = search_in_text(text, keywords)
                for match in matches:
                    if match_count >= limit:
                        break
                    snippet = create_snippet(text, match)
                    if snippet:
                        snippets.append(snippet)
                        match_count += 1
                if match_count >= limit:
                    break
        if snippets:
            copied_file = copy_file_to_download_folder(file_path)
            json_report[copied_file] = {
                "original_path": file_path,
                "snippets": snippets
            }
            print(f"{BOLD_TEXT}{PURPLE_TEXT}{file_path}{RESET_TEXT}:")
            for snippet in snippets:
                print(f"{snippet}\n")
    except Exception as e:
        debug_print(f"Skipping PDF file {file_path} due to error: {e}")
    return snippets


# Function to handle .docx files
def handle_docx_file(file_path, keywords, limit=10):
    """Read and search in .docx files."""
    snippets = []
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text
        match_count = 0  # Counter for number of matches
        matches = search_in_text(text, keywords)
        for match in matches:
            if match_count >= limit:
                break
            snippet = create_snippet(text, match)
            if snippet:
                snippets.append(snippet)
                match_count += 1
        if snippets:
            copied_file = copy_file_to_download_folder(file_path)
            json_report[copied_file] = {
                "original_path": file_path,
                "snippets": snippets
            }
            print(f"{BOLD_TEXT}{PURPLE_TEXT}{file_path}{RESET_TEXT}:")
            for snippet in snippets:
                print(f"{snippet}\n")
    except Exception as e:
        debug_print(f"Skipping .docx file {file_path} due to error: {e}")
    return snippets


# Function to handle .xlsx files
def handle_xlsx_file(file_path, keywords, limit=10):
    """Read and search in .xlsx files."""
    snippets = []
    try:
        wb = openpyxl.load_workbook(file_path)
        match_count = 0  # Counter for number of matches
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if match_count >= limit:
                        break
                    matches = search_in_text(str(cell), keywords)
                    for match in matches:
                        if match_count >= limit:
                            break
                        snippet = create_snippet(str(cell), match)
                        if snippet:
                            snippets.append(snippet)
                            match_count += 1
                if match_count >= limit:
                    break
            if match_count >= limit:
                break
        if snippets:
            copied_file = copy_file_to_download_folder(file_path)
            json_report[copied_file] = {
                "original_path": file_path,
                "snippets": snippets
            }
            print(f"{BOLD_TEXT}{PURPLE_TEXT}{file_path}{RESET_TEXT}:")
            for snippet in snippets:
                print(f"{snippet}\n")
    except Exception as e:
        debug_print(f"Skipping .xlsx file {file_path} due to error: {e}")
    return snippets


# Function to handle .pptx files
def handle_pptx_file(file_path, keywords, limit=10):
    """Read and search in .pptx files."""
    snippets = []
    try:
        prs = Presentation(file_path)
        match_count = 0  # Counter for number of matches
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if match_count >= limit:
                        break
                    matches = search_in_text(shape.text, keywords)
                    for match in matches:
                        if match_count >= limit:
                            break
                        snippet = create_snippet(shape.text, match)
                        if snippet:
                            snippets.append(snippet)
                            match_count += 1
                if match_count >= limit:
                    break
            if match_count >= limit:
                break
        if snippets:
            copied_file = copy_file_to_download_folder(file_path)
            json_report[copied_file] = {
                "original_path": file_path,
                "snippets": snippets
            }
            print(f"{BOLD_TEXT}{PURPLE_TEXT}{file_path}{RESET_TEXT}:")
            for snippet in snippets:
                print(f"{snippet}\n")
    except Exception as e:
        debug_print(f"Skipping .pptx file {file_path} due to error: {e}")
    return snippets


# Handle .odt files
def handle_odt_file(file_path, keywords, limit=10):
    """Read and search in .odt files."""
    snippets = []
    try:
        doc = load_odt(file_path)
        text = ""
        for paragraph in doc.getElementsByType(P):
            text += paragraph.text
        match_count = 0  # Counter for number of matches
        matches = search_in_text(text, keywords)
        for match in matches:
            if match_count >= limit:
                break
            snippet = create_snippet(text, match)
            if snippet:
                snippets.append(snippet)
                match_count += 1
        if snippets:
            copied_file = copy_file_to_download_folder(file_path)
            json_report[copied_file] = {
                "original_path": file_path,
                "snippets": snippets
            }
            print(f"{BOLD_TEXT}{PURPLE_TEXT}{file_path}{RESET_TEXT}:")
            for snippet in snippets:
                print(f"{snippet}\n")
    except Exception as e:
        debug_print(f"Skipping .odt file {file_path} due to error: {e}")
    return snippets

def handle_interesting_filename(file_path, keywords, limit):
    interesting_filenames.append(file_path)
    

# List of handler functions for each extension
file_handlers = {
    '.pdf': handle_pdf_file,
    '.docx': handle_docx_file,
    '.xlsx': handle_xlsx_file,
    '.pptx': handle_pptx_file,
    '.odt': handle_odt_file,
    '.wim': handle_interesting_filename,
    '.kdbx': handle_interesting_filename,
}


# Function to handle files based on their type
def handle_file(file_path, keywords, limit):
    """Determine file type and search for keywords."""
    extension = os.path.splitext(file_path)[1].lower()
    handler = file_handlers.get(extension)
    
    if handler:
        return handler(file_path, keywords, limit)
    else:
        debug_print(f"No handler for {file_path}")
    return None


# Main function to scan the directory with multi-threading
def scan_directory(directory, keywords, limit):
    """Scan the directory for files and search for keywords using multi-threading."""
    matched_files = []

    # Use ThreadPoolExecutor to handle files in parallel
    with ThreadPoolExecutor() as executor:
        futures = []
        
        for root, _, files in os.walk(directory):
            for file in files:
                file_path = os.path.join(root, file)
                futures.append(executor.submit(handle_file, file_path, keywords, limit))
        
        # Process results
        for future in as_completed(futures):
            result = future.result()
            if result:
                matched_files.append(result)

    return matched_files


def debug_print(message):
    """Print debug messages if in debug mode."""
    if debug:
        print(f"{RED_TEXT}{message}{RESET_TEXT}")



def analyze_occurrences_in_report(report_file_path, model_name="artifish/llama3.2-uncensored"):
    # Open and load the JSON report
    with open(report_file_path, "r", encoding="utf-8") as file:
        report_data = json.load(file)


    # Iterate over the entries in the report data (assuming the structure is similar to the example provided)
    for file_path, file_data in report_data.items():
        if 'snippets' in file_data:
            for snippet in file_data['snippets']:
                # Assuming the snippet is a text containing the occurrence we want to analyze
                occurrence_text = snippet  # Assuming text after the colon is the occurrence text

                # Format the occurrence text for the analysis
                formatted_text = f"{{{{{occurrence_text}}}}}"

                message_content = (
                    f"does the following text contain credentials? \nAnswer with the following format\n if no password : OK\n if only password: pass=the password\n if password and user user=password\nHere is the text:\n {formatted_text}"
                )

                # Send the request to the AI model
                response = ollama.chat(model=model_name, messages=[{"role": "user", "content": message_content}])

                result = response['message']['content'].strip()
                print(result)



# Get the directory of the script
script_directory = os.path.dirname(os.path.realpath(__file__))
default_keywords_file = os.path.join(script_directory, "keywords.txt")

# Main entry point for the script
if __name__ == "__main__":
    # Set up argument parsing
    parser = argparse.ArgumentParser(description="Search files for keywords")
    parser.add_argument("directory", help="Directory to scan for files")
    parser.add_argument("--keywords", help="File containing keywords", default=default_keywords_file)
    parser.add_argument("--debug", help="Enable debug mode", action="store_true")
    parser.add_argument("--limit", type=int, default=10, help="Limit number of keyword matches per file")
    parser.add_argument("--nollm", help="Disable llm", action="store_true")
    parser.add_argument("--json", default=None, help="Specify a json report, and treat it with llm")

    args = parser.parse_args()

    debug = args.debug

    if not args.json:
        # Load keywords
        keywords_file = args.keywords if os.path.isfile(args.keywords) else "keywords.txt"
        keywords = load_keywords(keywords_file)

        # Scan the directory for files containing the keywords
        matched_files = scan_directory(args.directory, keywords, args.limit)

        # Write the JSON report
        json_report_path = os.path.join(download_folder, "report.json")
        with open(json_report_path, "w") as json_file:
            json.dump(json_report, json_file, indent=4)

        print(f"\n{BOLD_TEXT}Interesting filenames:{RESET_TEXT}")
        for filename in interesting_filenames:
            print(filename)

        print(f"\nJSON report saved to {json_report_path}")
    
    if not args.nollm:
        if args.json:
            report_path = args.json
        else:
            report_path = json_report_path
        analyze_occurrences_in_report(report_path)

